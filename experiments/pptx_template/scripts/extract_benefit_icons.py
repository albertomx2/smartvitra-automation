from pathlib import Path

from pptx import Presentation

TEMPLATE = Path("experiments/pptx_template/input/template.pptx")

OUTPUT = Path("assets/presentation/icons/benefits")

OUTPUT.mkdir(
    parents=True,
    exist_ok=True,
)

prs = Presentation(str(TEMPLATE))

slide = prs.slides[6]

names = {
    "sv_s07_benefit_1_icon": "thermal",
    "sv_s07_benefit_2_icon": "acoustic",
    "sv_s07_benefit_3_icon": "energy",
    "sv_s07_benefit_4_icon": "home_value",
}

for shape in slide.shapes:
    if shape.name not in names:
        continue

    image = shape.image

    extension = image.ext

    output_path = OUTPUT / f"{names[shape.name]}.{extension}"

    output_path.write_bytes(image.blob)

    print(
        shape.name,
        "->",
        output_path,
    )
