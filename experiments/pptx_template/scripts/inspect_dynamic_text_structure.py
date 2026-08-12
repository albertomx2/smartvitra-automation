from pathlib import Path

from pptx import Presentation

TEMPLATE = Path("experiments/pptx_template/input/template.pptx")


TARGETS = {
    3: {
        "sv_s03_main_benefit",
        "sv_s03_main_benefit_secondary",
        "sv_s03_benefit_claim",
    },
    7: {
        "sv_s07_project_summary",
        "sv_s07_budget_block",
    },
}


prs = Presentation(TEMPLATE)


for slide_number, names in TARGETS.items():
    slide = prs.slides[slide_number - 1]

    print()
    print("#" * 80)
    print(f"SLIDE {slide_number}")
    print("#" * 80)

    for shape in slide.shapes:
        if shape.name not in names:
            continue

        print()
        print("=" * 80)
        print(shape.name)
        print("=" * 80)

        if not shape.has_text_frame:
            print("NO TEXT FRAME")
            continue

        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
            print()
            print(f"PARAGRAPH {paragraph_index}: " f"{paragraph.text!r}")

            for run_index, run in enumerate(paragraph.runs):
                font = run.font

                print(
                    "  "
                    f"RUN {run_index}: "
                    f"text={run.text!r} "
                    f"bold={font.bold} "
                    f"italic={font.italic} "
                    f"size={font.size}"
                )
