from pathlib import Path

from pptx import Presentation

TEMPLATE = Path("experiments/pptx_template/input/template.pptx")


RENAMES: dict[int, dict[str, str]] = {
    # ------------------------------------------------------------------
    # SLIDE 1
    # ------------------------------------------------------------------
    1: {
        "object 2": "sv_s01_cover_photo",
        "object 5": "sv_s01_intro_text",
        "object 7": "sv_s01_customer_name",
        "object 9": "sv_s01_address",
        "object 11": "sv_s01_proposal_number",
        "object 13": "sv_s01_date",
    },
    # ------------------------------------------------------------------
    # SLIDE 2
    # ------------------------------------------------------------------
    2: {
        "object 2": "sv_s02_problem_photo",
        "object 5": "sv_s02_issue_1",
        "object 7": "sv_s02_issue_2",
        "object 9": "sv_s02_issue_3",
        "object 11": "sv_s02_issue_4",
        "object 13": "sv_s02_issue_5",
        "object 15": "sv_s02_issue_6",
    },
    # ------------------------------------------------------------------
    # SLIDE 3
    #
    # object 2 = imagen grande de la derecha
    #
    # Los iconos van emparejados 1:1 con cada solución.
    # ------------------------------------------------------------------
    3: {
        "object 2": "sv_s03_generated_solution_image",
        "object 4": "sv_s03_solution_1_icon",
        "object 5": "sv_s03_solution_1",
        "object 6": "sv_s03_solution_2_icon",
        "object 7": "sv_s03_solution_2",
        "object 8": "sv_s03_solution_3_icon",
        "object 9": "sv_s03_solution_3",
        "object 10": "sv_s03_solution_4_icon",
        "object 11": "sv_s03_solution_4",
        "object 12": "sv_s03_solution_5_icon",
        "object 13": "sv_s03_solution_5",
        "object 14": "sv_s03_solution_6_icon",
        "object 15": "sv_s03_solution_6",
        # Este cuadro contiene:
        # "Beneficio principal"
        # +
        # el texto principal.
        "object 19": "sv_s03_main_benefit",
        "object 20": "sv_s03_main_benefit_secondary",
        "object 21": "sv_s03_benefit_claim",
    },
    # ------------------------------------------------------------------
    # SLIDE 5
    # ------------------------------------------------------------------
    5: {
        "object 3": "sv_s05_project_photo_1",
        "object 4": "sv_s05_project_photo_2",
        "object 5": "sv_s05_project_photo_3",
    },
    # ------------------------------------------------------------------
    # SLIDE 7
    #
    # Los bloques de texto son deliberadamente bloques completos.
    # Después modificaremos sus párrafos/runs internamente.
    # ------------------------------------------------------------------
    7: {
        "object 2": "sv_s07_generated_result_image",
        "object 4": "sv_s07_project_summary",
        "object 5": "sv_s07_budget_block",
    },
}


prs = Presentation(TEMPLATE)

renamed = 0


for slide_number, mapping in RENAMES.items():
    slide = prs.slides[slide_number - 1]

    shapes_by_name = {shape.name: shape for shape in slide.shapes}

    print()
    print("=" * 80)
    print(f"SLIDE {slide_number}")
    print("=" * 80)

    for old_name, new_name in mapping.items():
        shape = shapes_by_name.get(old_name)

        if shape is None:
            raise RuntimeError(f"Slide {slide_number}: " f"no existe {old_name!r}")

        shape.name = new_name

        renamed += 1

        print(f"{old_name:12} -> {new_name}")


prs.save(TEMPLATE)

print()
print("=" * 80)
print(f"{renamed} objetos renombrados")
print(TEMPLATE)
