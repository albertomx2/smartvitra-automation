from pathlib import Path

from pptx import Presentation

template_path = Path("experiments/pptx_template/input/template.pptx")

presentation = Presentation(template_path)

print(f"Slides: {len(presentation.slides)}")

for slide_index, slide in enumerate(
    presentation.slides,
    start=1,
):
    print()
    print("=" * 80)
    print(f"SLIDE {slide_index}")
    print("=" * 80)

    for shape in slide.shapes:
        has_text = getattr(
            shape,
            "has_text_frame",
            False,
        )

        print(
            f"name={shape.name!r} " f"type={shape.shape_type} " f"has_text={has_text}"
        )

        if has_text:
            text = shape.text.replace(
                "\n",
                " | ",
            )

            print(f"  text={text!r}")
